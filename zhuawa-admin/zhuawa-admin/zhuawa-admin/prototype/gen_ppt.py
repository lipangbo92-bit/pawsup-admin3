from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(16)

# Color scheme
PRIMARY = RGBColor(236, 72, 153)  # Pink
TEXT_COLOR = RGBColor(131, 24, 67)
GRAY_COLOR = RGBColor(107, 114, 128)

def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = GRAY_COLOR
        p2.alignment = PP_ALIGN.CENTER

def add_page_slide(prs, page_num, title, description):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Page number badge
    badge = slide.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(0.8), Inches(0.8))  # Oval
    badge.fill.solid()
    badge.fill.fore_color.rgb = PRIMARY
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].text = str(page_num)
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(1))
    tf = txBox.text_frame
    tf.paragraphs[0].text = f"页面 {page_num}: {title}"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = TEXT_COLOR
    
    # Description
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = description
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_COLOR
    p.space_after = Pt(10)
    
    # Phone frame placeholder
    frame = slide.shapes.add_shape(1, Inches(3.5), Inches(5), Inches(3), Inches(6))
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(253, 248, 252)
    frame.line.color.rgb = PRIMARY
    frame.line.width = Pt(2)
    tf = frame.text_frame
    tf.paragraphs[0].text = f"📱 {title}\n\n(见HTML原型)"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.color.rgb = RGBColor(156, 163, 175)

# Pages data
pages = [
    (1, "首页", "服务项目列表 + 技师展示\n- 店铺名称 + 位置\n- 金刚区：洗美护理、会员等入口\n- 技师卡片横向滑动\n- 热门服务列表"),
    (2, "技师详情", "技师详细信息\n- 大头像 + 名字 + 级别\n- 擅长服务标签\n- 评分 + 预约数\n- 立即预约按钮"),
    (3, "预约-选择服务", "选择服务项目\n- 步骤条显示进度\n- 服务列表：名称、价格、时长"),
    (4, "预约-选择技师", "选择技师\n- 根据已选服务和日期筛选\n- 技师卡片展示"),
    (5, "预约-选择时间", "选择预约时间\n- 日期选择器（横向日期条）\n- 时间段网格（30分钟单位）\n- 已预约时间显示不可选\n- 预约信息摘要"),
    (6, "选择宠物", "选择宠物\n- 已有宠物列表\n- 新增宠物入口"),
    (7, "预约确认", "确认订单信息\n- 服务项目、技师、宠物\n- 日期时间、时长\n- 总金额\n- 取消政策提示"),
    (8, "支付", "微信支付\n- 订单金额展示\n- 微信支付按钮"),
    (9, "预约成功", "预约成功提示\n- 成功图标动画\n- 预约详情摘要\n- 查看预约 / 返回首页"),
    (10, "预约记录", "我的预约列表\n- Tab筛选：全部/待消费/已完成/已取消\n- 预约卡片列表"),
    (11, "预约详情", "预约详细信息\n- 状态进度条\n- 预约信息\n- 取消预约 / 联系技师按钮"),
]

# Title slide
add_title_slide(prs, "🐾 宠物店预约小程序", "顾客端原型图 - 交互流程")

# Page slides
for page_num, title, description in pages:
    add_page_slide(prs, page_num, title, description)

# Save
output_path = os.path.expanduser("~/Desktop/pet_store_prototype.pptx")
prs.save(output_path)
print(f"PPT已生成: {output_path}")
